import pptx


def get_text_from_pptx(filename: str) -> dict[int, list[str]]:
    result = {}
    slides = pptx.Presentation(filename).slides
    for slide in slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_number = slides.index(slide)
                if slide_number not in result:
                    result[slide_number] = []
                result[slide_number].append(shape.text)
    return result


if __name__ == "__main__":
    texts = get_text_from_pptx('pres.pptx')
    for index in texts:
        print(f"{index} {texts[index]}")
