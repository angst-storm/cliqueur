import asyncio
import json
import logging
import re
from tempfile import NamedTemporaryFile
import pptx
from pptx import Presentation
import io

import gigachat_handler
import gigachat_pic
from s3 import s3_resource, s3_client, BUCKET_NAME
import boto3
import uuid
import os

import aspose.slides as slides
from aspose.slides.export import HtmlOptions, SaveFormat
from fastapi import WebSocket
from fastapi import WebSocketDisconnect
from json import loads

PRESENTATION_LINK_BASE = os.getenv("PRESENTATION_LINK_BASE")
PROCESS_IMAGES = os.getenv("PROCESS_IMAGES", "false") == "true"

MIN_PROBABILITY = 0.4
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

slides_queue = asyncio.Queue()
pres_status = {'isContextMode': False, 'isKeywordMode': False, 'currentSlide': 0}


class PresentationConverter:
    def __init__(self, file_data: bytes):
        temp_pptx = NamedTemporaryFile(delete=False, suffix=".pptx")
        temp_pptx.write(file_data)
        self.temp_pptx_path = temp_pptx.name

    def convert_to_html(self) -> str:
        try:
            with slides.Presentation(self.temp_pptx_path) as pres:
                options = HtmlOptions()
                output_path = self.temp_pptx_path.replace(".pptx", ".html")
                pres.save(output_path, SaveFormat.HTML, options)

            with open(output_path, "r", encoding="utf-8") as f:
                html_content = f.read()

            html_content = re.sub(
                r"<tspan[^>]*>Evaluation only\.</tspan>", "", html_content
            )
            html_content = re.sub(
                r"<tspan[^>]*>Created with Aspose\.Slides[^<]*</tspan>",
                "",
                html_content,
            )
            html_content = re.sub(
                r"<tspan[^>]*>Copyright \d{4}-\d{4}Aspose Pty Ltd\.</tspan>",
                "",
                html_content,
            )
            html_content = re.sub(
                r'<div\s+class="slideTitle">.*?</div>', "", html_content
            )
        except Exception as e:
            logger.info(f"HTML convertation error: {e}")
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

        return html_content

    def __del__(self):
        os.unlink(self.temp_pptx_path)
        logger.info("Pres converter deleted")


async def process_presentation(websocket: WebSocket):
    await websocket.accept()
    logger.info("Pres WebSocket подключен")

    try:
        pptx_data = await websocket.receive_bytes()
        logger.info("Получен файл %s байт", len(pptx_data))

        pres_converter = PresentationConverter(pptx_data)
        html = pres_converter.convert_to_html()
        notes_map = extract_bracketed_notes(pptx_data)
        logger.info("Converted to HTML")

        pres_id = uuid.uuid4()

        save_s3(pres_id, html, pptx_data)
        save_s3_keywords(pres_id, notes_map)
        slides_text = extract_text(pres_id)
        logger.info("Text extracted")

        if PROCESS_IMAGES:
            giga_pic = gigachat_pic.GigachatPicHandler()
            processed_images = giga_pic.process_aspose(pres_converter.temp_pptx_path)
            slides_text = combine_text_and_images(slides_text, processed_images)
            logger.info("Images processed")

        giga_proc = gigachat_handler.GigachatPresHandler()  # todo это надо в очредь какую-нибудь
        preprocess_text = giga_proc.process_presentation(slides_text, pres_id)

        save_s3_preprocess(pres_id, preprocess_text)

        link = f"{PRESENTATION_LINK_BASE}/{pres_id}"
        await websocket.send_text(link)
        logger.info("Отправлена ссылка %s", link)

    except WebSocketDisconnect:
        logger.info("Клиент pres отключился")
    except Exception as e:
        logger.error("Ошибка обработки: %s", e)
        await websocket.send_text(f"Ошибка конвертации: {str(e)}")

    finally:
        await websocket.close()


def combine_text_and_images(text_dict: dict, images: dict[int, str]):
    combined = text_dict.copy()
    for i in images:
        combined[i]["images"] = images[i]

    return combined


def extract_bracketed_notes(pptx_data: bytes) -> dict[int, list[str]]:
    prs = Presentation(io.BytesIO(pptx_data))
    pattern = re.compile(r'\[([^]]+)]')
    notes_map: dict[int, list[str]] = {}

    for idx, slide in enumerate(prs.slides):
        if not slide.has_notes_slide:
            continue

        notes: list[str] = []
        for shape in slide.notes_slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text or ""
            notes.extend(pattern.findall(text))

        if notes:
            notes_map[idx] = notes

    return notes_map


def extract_text(pres_id: str) -> dict[int, dict]:
    data = io.BytesIO()
    s3_client.download_fileobj(
        Bucket=BUCKET_NAME, Key=f"{pres_id}/file.pptx", Fileobj=data
    )
    pres = pptx.Presentation(data)
    text = {}
    for idx, slide in enumerate(pres.slides):
        text[idx] = {"text": [], "images": None}
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text[idx]["text"].append(shape.text)

    return text


async def send_slide_number(websocket: WebSocket):
    logger.info("Slides Send WebSocket подключен")
    try:
        while True:
            slides_probs = await slides_queue.get()
            if not slides_probs:
                logger.info("Slides: skipping empty dict")
                continue

            slide = max(slides_probs.items(), key=lambda x: x[1])
            logger.info(f"Top slide is {slide}")
            if slide[1] >= MIN_PROBABILITY:
                logger.info(f"Confident enough. Sending number {slide[0]}")
                await websocket.send_text(str(slide[0]))
    except WebSocketDisconnect:
        logger.info("Клиент slides send отключился")
    except Exception as e:
        logger.error(f"Ошибка отправки слайда: {str(e)}")
    finally:
        await websocket.close()


async def get_front_status(websocket: WebSocket):
    global pres_status
    logger.info("Slides Receive WebSocket подключен")
    try:
        while True:
            status = await websocket.receive_text()
            pres_status = loads(status)
            logger.info(pres_status)
    except WebSocketDisconnect:
        logger.info("Клиент slides receive отключился")
    except Exception as e:
        logger.error(f"Ошибка отправки слайда: {str(e)}")


def save_s3(pres_id: str, html: str, pptx_data: bytes):
    html_object = s3_resource.Object(BUCKET_NAME, f"{pres_id}/index.html")
    html_object.put(Body=html)

    pptx_object = s3_resource.Object(BUCKET_NAME, f"{pres_id}/file.pptx")
    pptx_object.put(Body=pptx_data)

    logger.info("Презентация %s успешно сохранена в S3", pres_id)


def save_s3_preprocess(pres_id: str, preprocess_text: str):
    preprocess_object = s3_resource.Object(BUCKET_NAME, f"{pres_id}/preprocess.json")
    preprocess_object.put(Body=preprocess_text)

    logger.info("Предобработка %s успешно сохранена в S3", pres_id)


def save_s3_keywords(pres_id: str, keywords_map: dict[int, list[str]]):
    body = json.dumps(keywords_map, ensure_ascii=False)
    keywords_object = s3_resource.Object(BUCKET_NAME, f"{pres_id}/keywords.json")
    keywords_object.put(Body=body)

    logger.info("Ключевые слова %s успешно сохранены в S3", pres_id)


def load_s3_keywords(pres_id: str) -> dict[int, list[str]]:
    obj = s3_client.get_object(Bucket=BUCKET_NAME, Key=f"{pres_id}/keywords.json")
    raw = json.loads(obj["Body"].read().decode("utf-8"))
    return {int(k): v for k, v in raw.items()}
